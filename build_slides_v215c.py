#!/usr/bin/env python3
"""Build a ~30-min general-audience PPTX from manuscript v2.15c.

Usage: python build_slides.py   (from repo root)
Output: HLA_Registry_Size_CMIO_v2.15c_slides.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
import os

# Frozen 1e-3 snapshot: live analysis/figures now holds the 1e-6 re-run
# (e.g. em_convergence.png rises to 8.7e7, contradicting the v2.15 caption).
FIG = "analysis/snapshot_1e-3/figures"
OUT = "HLA_Registry_Size_CMIO_v2.15c_slides.pptx"

# neutral palette
INK = RGBColor(0x1F, 0x29, 0x37)      # near-black text
ACCENT = RGBColor(0x1F, 0x3A, 0x5F)   # deep navy
MUTED = RGBColor(0x6B, 0x72, 0x80)    # gray
HEAD_BG = RGBColor(0xEE, 0xF1, 0xF5)  # table header fill
RED = RGBColor(0xB4, 0x23, 0x18)

SW, SH = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def txt(s, l, t, w, h, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, font="Calibri"):
    box = s.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    f = r.font
    f.size, f.bold, f.name = Pt(size), bold, font
    f.color.rgb = color
    return box


def title(s, text, sub=None):
    txt(s, Inches(0.55), Inches(0.32), Inches(12.2), Inches(0.9), text, size=30, bold=True, color=ACCENT)
    ln = s.shapes.add_shape(1, Inches(0.6), Inches(1.12), Inches(12.1), Pt(2.2))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT; ln.line.fill.background()
    if sub:
        txt(s, Inches(0.55), Inches(1.2), Inches(12.2), Inches(0.5), sub, size=15, color=MUTED)


def bullets(s, items, l=Inches(0.7), t=Inches(1.6), w=Inches(11.9), h=Inches(5.4), size=20, gap=8):
    box = s.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for it in items:
        lvl, text, kw = 0, it, {}
        if isinstance(it, tuple):
            if isinstance(it[0], str):           # (text, kw)
                text, kw = it[0], it[1]
            else:                                 # (lvl, text[, kw])
                lvl, text = it[0], it[1]
                kw = it[2] if len(it) > 2 else {}
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        p.space_after = Pt(gap)
        r = p.add_run()
        r.text = ("• " if lvl == 0 else "– ") + text
        f = r.font
        f.size = Pt(kw.get("size", size - 3 * lvl))
        f.bold = kw.get("bold", False)
        f.name = "Calibri"
        f.color.rgb = kw.get("color", INK)
    return box


def pic(s, path, l, t, max_w, max_h):
    iw, ih = Image.open(path).size
    scale = min(max_w / iw, max_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    # center within the box
    return s.shapes.add_picture(path, int(l + (max_w - w) / 2), int(t + (max_h - h) / 2), w, h)


def table(s, rows, l, t, w, h, size=13, header=True, col_w=None, bold_cols=(), red_cells=()):
    nr, nc = len(rows), len(rows[0])
    shp = s.shapes.add_table(nr, nc, l, t, w, h)
    tb = shp.table
    if col_w:
        total = sum(col_w)
        for i, cw in enumerate(col_w):
            tb.columns[i].width = int(w * cw / total)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            c = tb.cell(ri, ci)
            c.margin_top = c.margin_bottom = Pt(2)
            tf = c.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run(); r.text = str(val)
            f = r.font
            f.size = Pt(size); f.name = "Calibri"
            hdr = header and ri == 0
            f.bold = hdr or ci in bold_cols
            f.color.rgb = RED if (ri, ci) in red_cells else (ACCENT if hdr else INK)
            c.fill.solid()
            c.fill.fore_color.rgb = HEAD_BG if hdr else RGBColor(0xFF, 0xFF, 0xFF)
    return shp


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def caption(s, text, t=Inches(6.85)):
    txt(s, Inches(0.7), t, Inches(11.9), Inches(0.5), text, size=12, color=MUTED)


# ---------------------------------------------------------------- 1 title
s = slide()
txt(s, Inches(1.0), Inches(2.2), Inches(11.3), Inches(1.8),
    "How Many Donors Does Singapore Need?", size=40, bold=True, color=ACCENT)
txt(s, Inches(1.0), Inches(3.5), Inches(11.3), Inches(1.0),
    "Modelling bone marrow registry size for Singapore's multiethnic (CMIO) population", size=22, color=INK)
txt(s, Inches(1.0), Inches(5.2), Inches(11.3), Inches(1.2),
    "Alvin Ng Yu-Jin\nNational University Hospital / Singapore Cord Blood Bank\nManuscript v2.15c · 2026", size=16, color=MUTED)
pic(s, "assets/cartoons/sg_island.png", Inches(8.5), Inches(4.0), Inches(4.3), Inches(3.1))
notes(s, "Good morning everyone. Today I'll talk about a deceptively simple question: how many volunteer bone marrow donors does Singapore actually need? Behind that question is a serious equity problem — patients from different ethnic backgrounds have very different chances of finding a life-saving donor. This work turns Singapore's largest-ever HLA dataset into concrete, defensible recruitment targets for each community. I'll keep the mathematics light and focus on the intuition, the pictures, and what the numbers mean for policy. The talk is about 30 minutes, and I'm happy to take questions at the end. [~1 min]")

# ---------------------------------------------------------------- 2 motivation
s = slide()
title(s, "Why this matters")
bullets(s, [
    "A stem cell transplant (HSCT) can cure many blood cancers and marrow disorders",
    "Most patients have no matched donor in their own family",
    "Their best hope: an unrelated volunteer from a donor registry",
    "The donor's and patient's HLA tissue types must match closely — mismatches raise the risk of rejection and graft-versus-host disease",
    ("A registry is only useful if it is large enough — and diverse enough — to contain a match", {"bold": True}),
], w=Inches(7.5), size=21)
pic(s, "assets/cartoons/patient_donor.png", Inches(8.4), Inches(2.0), Inches(4.5), Inches(3.6))
notes(s, "Let me start with the motivation. Haematopoietic stem cell transplantation — a bone marrow transplant — is a potentially curative treatment for leukaemias and other blood disorders. But it only works if we can find a donor whose tissue type, the HLA type, matches the patient's closely. Most patients — around 70% — don't have a matched sibling, so they depend on registries of unrelated volunteers. High-resolution HLA matching is strongly associated with survival. So the core question for any country running a registry is: how many donors do we need to sign up so that, when a patient walks through the door, a match is actually there? [~1.5 min]")

# ---------------------------------------------------------------- 3 the equity gap
s = slide()
title(s, "The problem is harder for non-European patients")
bullets(s, [
    "Worldwide registries are dominated by donors of European descent",
    "HLA types common in East Asian, South Asian, or mixed-ancestry patients are often rare or absent in those registries",
    "Singapore classifies its population as Chinese, Malay, Indian, Others — 'CMIO'",
    "Ng et al. (2022) mapped Singapore's HLA landscape across 59,186 donors — but did not translate it into registry size targets",
    ("This study asks: exactly how large must a local registry be, for each community?", {"bold": True, "color": ACCENT}),
], w=Inches(7.6), size=21)
pic(s, "assets/cartoons/crowd_gap.png", Inches(8.6), Inches(1.7), Inches(4.3), Inches(4.3))
notes(s, "Finding a donor is not equally hard for everyone. Global registries hold tens of millions of donors, but they are dominated by people of European descent. Tissue types that are common here in Singapore — among Chinese, Malay, Indian, or mixed-ancestry patients — may simply not exist in those registries. Singapore uses the CMIO classification: Chinese, Malay, Indian and Others. In 2022, our group published the HLA landscape of Singapore's donor pool — allele and haplotype frequencies for each group, from nearly 60,000 donors. That was the foundation. What was missing was the translation into planning numbers: how many donors, per community, for a given chance of finding a match. That's what this study delivers. [~1.5 min]")

# ---------------------------------------------------------------- 4 HLA matching in plain terms
s = slide()
title(s, "HLA matching in plain terms")
bullets(s, [
    "HLA genes are the immune system's 'ID badge' — the transplant works best when donor and patient badges match",
    "We look at 5 genes (A, B, C, DRB1, DQB1); everyone carries 2 copies of each → 10 slots",
    ("'10/10 match' = all ten slots identical (the strict clinical standard)", {"bold": True}),
    "'8/8 match' = the same, ignoring the DQB1 gene (4 genes × 2 copies)",
    "'9/10' = one slot mismatched — sometimes clinically acceptable",
    "Each gene comes in hundreds of versions, so exact matches between strangers are rare",
], w=Inches(7.4), size=20)
pic(s, "assets/cartoons/id_badges.png", Inches(8.3), Inches(1.8), Inches(4.6), Inches(4.1))
notes(s, "A quick plain-language primer, because the rest of the talk depends on it. HLA genes are like an ID badge your cells show to the immune system. If a transplanted cell's badge looks wrong, the immune system attacks — in both directions. We type five HLA genes: A, B, C, DRB1 and DQB1. You inherit two copies of each — one from each parent — giving ten 'slots'. A 10-out-of-10 match means all ten slots are identical: that's the strict standard. 8-out-of-8 ignores the DQB1 gene. And 9-out-of-10 means one slot is mismatched, which is sometimes acceptable clinically. The catch: each of these genes comes in hundreds of versions, so two strangers matching on all ten slots is genuinely rare — which is why registries need to be big. [~2 min]")

# ---------------------------------------------------------------- 5 four questions
s = slide()
title(s, "Four questions this study answers")
bullets(s, [
    ("1.  How many same-ethnicity donors are needed for 75–95% coverage, at 10/10 and 8/8 matching — and how confident are we?", {"bold": True}),
    "2.  Can a shared, Chinese-dominated pool serve Malay, Indian and Others patients (cross-ethnic matching)?",
    "3.  How much smaller could registries be if clinicians accepted partial matches (9/10 or 8/10)?",
    "4.  Is the 'Others' group really one population — or several, needing separate strategies?",
], w=Inches(8.3), size=21, gap=14)
pic(s, "assets/cartoons/questions.png", Inches(9.1), Inches(1.8), Inches(3.8), Inches(4.2))
notes(s, "The study answers four questions. First, the headline: how many same-ethnicity donors does each CMIO group need to give 75 to 95 per cent of its patients a match, at both match standards — and with honest error bars. Second, can we avoid ethnicity-specific recruitment altogether by relying on one big shared pool, which in Singapore would be predominantly Chinese? Third, what do we gain if clinicians relax the matching standard by one allele? And fourth, a question nobody had asked quantitatively: the 'Others' box on the census form contains Eurasians, Europeans, Filipinos, and many mixed backgrounds — is it legitimate to model them as one population? Keep these four in mind; the results section answers them one by one. [~1 min]")

# ---------------------------------------------------------------- 6 dataset
s = slide()
title(s, "The data: Singapore's donor pool")
bullets(s, [
    ("59,186 donors and cord blood units, accrued 2005–2020", {"bold": True}),
    "Sources: Bone Marrow Donor Programme (BMDP) + Singapore Cord Blood Bank (SCBB) — same cohort as Ng et al. 2022",
    "High-resolution typing at 5 genes: HLA-A, B, C, DRB1, DQB1",
    "Independent validation set: 564 real patient–donor pairs from the Health Sciences Authority (HSA)",
], size=20, t=Inches(1.5), h=Inches(2.6))
table(s, [
    ["Group", "Donors with full 5-gene typing", "Share"],
    ["Chinese", "44,400", "75.0%"],
    ["Malay", "5,578", "9.4%"],
    ["Indian", "5,490", "9.3%"],
    ["Others", "3,767", "6.4%"],
], Inches(0.9), Inches(4.15), Inches(6.5), Inches(2.6), size=15, col_w=[2, 3, 1.5])
pic(s, "assets/cartoons/lab_data.png", Inches(8.2), Inches(4.1), Inches(4.4), Inches(2.9))
notes(s, "The data. We used HLA typing from 59,186 donors and cord blood units collected between 2005 and 2020 by the Bone Marrow Donor Programme and the Singapore Cord Blood Bank — the largest local HLA dataset ever assembled. Everyone is typed at high resolution across the five genes I mentioned. This is exactly the cohort of the 2022 paper — same sources, same inclusion criteria — because this work is a direct follow-up on the identical dataset. Note the composition in the table: about three-quarters of fully-typed donors are Chinese, with roughly 5,600 Malay, 5,500 Indian and 3,800 Others. That imbalance matters twice: it drives the cross-ethnic result later, and it means our statistical confidence is highest for the Chinese group. We also secured an independent test set — 564 real patient–donor pairs from the Health Sciences Authority — which we use to check the model against reality. [~1.5 min]")

# ---------------------------------------------------------------- 7 pipeline figure
s = slide()
title(s, "How the analysis works — end to end")
pic(s, f"{FIG}/pipeline_flowchart.png", Inches(0.7), Inches(1.35), Inches(11.9), Inches(5.35))
caption(s, "Figure 1. Analysis pipeline: blue = input data, green = core model, yellow = optimisation, red = uncertainty quantification.")
notes(s, "Here's the whole pipeline on one slide — this is the map for the next three slides. We start on the left with the raw typing data in blue. The green core: we first reconstruct haplotypes — I'll explain what those are in a moment — then combine them into genotype frequencies, and feed those into a coverage model that predicts, for a registry of any given size, what fraction of patients would find a match. The yellow step searches for the smallest registry that hits a target coverage — say 95 per cent. And the red step repeats the whole calculation a thousand times on statistically perturbed data to put honest error bars on every number. Three ideas — haplotypes, coverage, and bootstrap — and that's the entire method. [~1.5 min]")

# ---------------------------------------------------------------- 8 haplotypes / LD
s = slide()
title(s, "Key idea 1: EM haplotype phasing — genes travel together in blocks")
bullets(s, [
    "HLA genes sit side-by-side on chromosome 6 and are inherited as blocks called haplotypes",
    "Neighbouring genes are strongly linked: DRB1–DQB1 and B–C versions co-travel ≥94–95% of the time in every CMIO group",
    "Treating genes as independent (multiplying frequencies) badly underestimates how common the frequent combinations are",
    ("We used an expectation–maximisation (EM) algorithm to statistically reconstruct each group's haplotype frequencies — validated against the international GENE[RATE] database", {"bold": True}),
], h=Inches(3.5), size=20, gap=12)
pic(s, "assets/cartoons/gene_train.png", Inches(0.9), Inches(5.2), Inches(11.5), Inches(1.9))
notes(s, "Key idea number one. The five HLA genes sit next to each other on chromosome 6, so they're inherited together in blocks — a block is called a haplotype. And the linkage is strong: in every CMIO group, particular versions of DRB1 and DQB1 travel together at least 94 per cent of the time, similarly for B and C. Why does this matter? If you naively assume the genes are independent and just multiply frequencies, you dramatically underestimate how common the popular combinations are — and you'd overestimate the registry you need in the wrong places. So a key methodological choice was to reconstruct full five-gene haplotype frequencies with an expectation–maximisation algorithm — a standard statistical technique — and we validated the output against the international GENE[RATE] reference pipeline. [~1.5 min]")

# ---------------------------------------------------------------- 9 coverage intuition
s = slide()
title(s, "Key idea 2: the coverage model — from frequencies to a registry size")
bullets(s, [
    "For each patient tissue type, ask: in a registry of N random donors, what is the chance at least one donor carries the same type?",
    "Common types are matched by small registries; rare types need very large ones",
    ("'Coverage' = the expected % of patients who find ≥1 match in a registry of size N", {"bold": True}),
    "We then search for N* — the smallest registry achieving a target coverage (75%, 85%, 90%, 95%)",
    "Same framework pioneered for the US registry (Beatty 1995; Gragert NEJM 2014)",
], w=Inches(7.7), size=20, gap=12)
pic(s, "assets/cartoons/umbrella.png", Inches(8.7), Inches(1.8), Inches(4.1), Inches(4.2))
notes(s, "Key idea two — and I'll do this without equations. Take one patient with a particular tissue type. If that type occurs in, say, 1 in 1,000 people, then a registry of N random donors contains at least one match with a probability you can compute directly — it grows with N. Do this for every possible tissue type, weight by how many patients carry each type, and you get 'coverage': the expected percentage of patients who find at least one matched donor in a registry of size N. Coverage rises steeply at first, then flattens — the common types get covered early, the rare ones very late. Finally we run a search for N-star: the smallest registry that reaches a chosen coverage target. This is the same framework used to plan the US national registry, published in the New England Journal in 2014. [~1.5 min]")

# ---------------------------------------------------------------- 10 bootstrap
s = slide()
title(s, "Key idea 3: bootstrap uncertainty — honest error bars")
bullets(s, [
    "Haplotype frequencies are estimates from a finite sample — they carry statistical noise",
    "We re-ran the entire calculation 1,000 times, each time re-drawing plausible frequencies (Dirichlet bootstrap)",
    ("Reported value = the median of 1,000 answers; the middle 95% of answers = the confidence interval", {"bold": True}),
    ("Intervals are comparable across groups (±800–1,100 donors at 95%): the 5,000-person EM limit "
     "equalises the sample the frequencies come from", {"color": RED}),
    "The CIs capture frequency-sampling noise only — a stated lower bound on total uncertainty",
], w=Inches(7.7), size=20, gap=12)
pic(s, "assets/cartoons/bootstrap_icon.png", Inches(8.8), Inches(2.0), Inches(3.9), Inches(3.9))
notes(s, "Key idea three: uncertainty. Our haplotype frequencies come from a finite sample, so they're estimates, not truth. If we'd happened to recruit a slightly different 59,000 donors, the frequencies — and therefore the registry targets — would shift. To quantify that, we used a bootstrap: we re-drew statistically plausible frequency sets 1,000 times and recomputed the registry size each time. We report the median of those 1,000 answers, and the middle 95 per cent gives a confidence interval. As you'd expect, the group with the most data — Chinese — has tight intervals, plus or minus about 200 donors; the smaller groups are plus or minus one to two thousand. One honesty note: these intervals capture frequency-sampling noise only, so they're a lower bound on total uncertainty — we say so explicitly in the paper. [~1.5 min]")

# ---------------------------------------------------------------- 11 Table 1
s = slide()
title(s, "Result 1 — donors needed for a 10/10 match (same ethnicity)",
      sub="Coverage is CONDITIONAL: of patients whose both haplotypes clear the 0.1% frequency floor. See the next slide.")
table(s, [
    ["Ethnicity", "75%", "85%", "90%", "95% coverage  (95% CI)"],
    ["Chinese", "7,292", "14,555", "22,534", "41,183  (40,184–42,153)"],
    ["Malay", "6,169", "12,712", "20,521", "39,831  (38,680–40,935)"],
    ["Indian", "9,491", "17,706", "26,038", "43,785  (42,762–44,547)"],
    ["Others (pooled)*", "6,864", "12,487", "18,341", "31,129  (30,420–32,001)"],
    ["Combined pooled registry†", "26,210", "60,875", "107,139", "243,849"],
    ["Signed-up target (40% attrition)‡", "10,282–15,818", "20,812–29,510", "30,568–43,397", "51,882–72,975"],
], Inches(0.7), Inches(1.5), Inches(11.9), Inches(3.9), size=14,
      col_w=[2.6, 1.3, 1.3, 1.3, 2.9], bold_cols=(4,),
      # every per-group figure moved when n_eff was matched to the EM sample
      red_cells=tuple((r, c) for r in (1, 2, 3, 4, 6) for c in range(1, 5)) + ((4, 0),))
bullets(s, [
    ("≈ 40,000–44,000 same-ethnicity donors per major group for 95% CONDITIONAL coverage", {"bold": True, "color": ACCENT}),
    "* Pooled Others figure is a statistical artefact — see the sub-group analysis later",
    "† / ‡ The combined pooled registry is a mathematical convenience, not a policy target; signed-up targets add ~1.67× for real-world donor attrition",
], t=Inches(5.55), size=16, gap=6)
notes(s, "Here is the headline result — Table 1 of the paper. Read the right-hand column: to give 95 per cent of patients a full 10-out-of-10 match from their own community — and hold on to the word 'conditional' here, I will unpack it on the very next slide — Singapore needs roughly 41,000 Chinese donors, 40,000 Malay donors, and 44,000 Indian donors, with the confidence intervals shown in brackets. So the memorable number is: about forty to forty-four thousand donors per major community. Two footnotes matter. The pooled 'Others' figure of 31,000 is misleading — I'll show you why in a few slides; don't use it for planning. And these are biologically matched donors: in the real world 30 to 50 per cent of registered donors are unreachable or deferred when called, so the bottom row grosses the targets up by about 1.67 times to actual sign-up numbers. [~2 min]")

# ------------------------------------------------- 11b the floor caveat (v2.15c Limitations)
s = slide()
title(s, "Reality check \u2014 what those registries actually deliver",
      sub="The 95% above is CONDITIONAL. Restoring the denominator gives coverage of ALL patients.")
table(s, [
    ["Group", "Donors typed\nat 5 loci", "Mass retained\nby 0.1% floor",
     "Model, all patients:\n95% \u00d7 mass retained", "Observed 10/10\nmatch rate"],
    ["Chinese", "44,400", "51.7%", "0.95 \u00d7 51.7% = 49.1%", "41.8%"],
    ["Malay", "5,578", "52.9%", "0.95 \u00d7 52.9% = 50.2%", "25.3%"],
    ["Indian", "5,490", "40.2%", "0.95 \u00d7 40.2% = 38.2%", "19.7%"],
    ["Others", "3,767", "35.6%", "0.95 \u00d7 35.6% = 33.8%", "18.1%"],
], Inches(0.7), Inches(1.95), Inches(11.9), Inches(2.5), size=14,
      col_w=[1.4, 1.8, 2.0, 3.2, 2.0], bold_cols=(4,))
bullets(s, [
    ("Singapore already holds 44,400 Chinese donors \u2014 essentially the modelled 95% target \u2014 "
     "yet the realised 10/10 match rate is 41.8%, not 95%", {"bold": True, "color": RED, "size": 18}),
    ("Only the Chinese row is like-for-like; the other registries are 5\u201310\u00d7 smaller than their own targets",
     {"size": 16}),
    ("The 0.1% frequency floor discards 47\u201364% of haplotype mass \u2014 the model covers 95% of what "
     "remains, which is ~49% of everyone", {"size": 16}),
], t=Inches(4.7), size=17, gap=8)
notes(s, "This is the most important caveat in the talk. The 95 per cent on the last slide is a conditional number. We applied a 0.1 per cent frequency floor to keep the mathematics tractable, and that floor throws away between 47 and 64 per cent of the haplotype mass. The model then covers 95 per cent of the patients who remain \u2014 not 95 per cent of everyone. Restore the denominator and the Chinese figure becomes 49 per cent. Now here is the test. Singapore already has 44,400 Chinese donors, essentially the size the model says we need. So what match rate do we actually observe? We can count it directly, with no model at all: how many donors have an exact 10-out-of-10 genotype twin somewhere else in the registry? The answer is 41.8 per cent. Against a model prediction of 49 per cent that is good agreement; against the headline 95 per cent it is a refutation. So please read the previous table as a conditional statement. The honest summary is that a registry of this size gives roughly two patients in five a full match \u2014 and the remaining three fifths are why the rest of this talk matters. [~2.5 min]")

# ---------------------------------------------------------------- 12 Figure 2 CI plot
s = slide()
title(s, "The same numbers, with their error bars")
pic(s, f"{FIG}/registry_ci_plot.png", Inches(1.2), Inches(1.35), Inches(10.9), Inches(5.35))
caption(s, "Figure 2. Bootstrap 95% confidence intervals for the minimum registry size at 95% coverage (10/10 matching), by CMIO group.")
notes(s, "This figure shows the 95-per-cent targets graphically, with their bootstrap confidence intervals as error bars. Two things to take away. First, the three major groups — Chinese, Malay, Indian — all land in the same band, roughly 40 to 44 thousand: haplotype diversity differs between the groups, but the order of magnitude of the task is the same for each community. Second, look at the widths: the Chinese interval is barely visible because we had 46,000 samples; Malay and Indian intervals are wider, reflecting five to six thousand samples each. The message for planners: the targets are statistically stable — the error bars are small relative to the numbers themselves — so these are actionable figures, not rough guesses. [~1 min]")

# ---------------------------------------------------------------- 13 long tail
s = slide()
title(s, "Why the last 5% of patients are so expensive")
pic(s, f"{FIG}/diplotype_longtail.png", Inches(0.7), Inches(1.35), Inches(7.6), Inches(5.3))
bullets(s, [
    ("Coverage targets get sharply more expensive past 90%", {"bold": True}),
    "Most patients carry common HLA combinations — covered by ~20,000 donors",
    "The final 5–10% carry rare combinations from a very long tail",
    "Covering them requires disproportionately many extra donors",
    "This is why 90% → 95% roughly doubles the target",
], l=Inches(8.5), t=Inches(1.6), w=Inches(4.4), size=17, gap=10)
caption(s, "Long tail of diplotype (tissue-type pair) frequencies: a few common types, thousands of rare ones.")
notes(s, "Why do the numbers explode as we push the coverage target up? This chart shows the frequency distribution of tissue-type combinations: a handful of very common types on the left, then an extremely long tail of rare ones. The first 90 per cent of patients carry the common combinations — you can serve them with about 20,000 donors. But the last five to ten per cent of patients each carry a combination that might occur once in tens of thousands of people. Every extra percentage point of coverage past 90 buys matches for progressively rarer patients, so the required registry roughly doubles between 90 and 95 per cent coverage. We still recommend the 95 per cent standard — the patients in that tail are disproportionately from minority and mixed-ancestry backgrounds, so stopping at 90 per cent quietly writes them off. [~1.5 min]")

# ---------------------------------------------------------------- 14 Table 2, 8/8
s = slide()
title(s, "Result 2 — dropping DQB1 (8/8) barely helps")
table(s, [
    ["Ethnicity", "75%", "85%", "90%", "95% coverage  (95% CI)"],
    ["Chinese", "7,188", "14,327", "22,170", "40,490  (39,534–41,528)"],
    ["Malay", "5,635", "11,548", "18,563", "35,999  (34,976–37,067)"],
    ["Indian", "9,169", "17,138", "25,256", "42,590  (41,628–43,391)"],
    ["Others (pooled)*", "6,454", "11,913", "17,704", "30,489  (29,790–31,271)"],
    ["Combined pooled registry†", "25,082", "57,847", "101,369", "229,800"],
    ["Signed-up target (40% attrition)‡", "9,392–15,282", "19,247–28,563", "29,507–42,093", "50,815–70,983"],
], Inches(0.7), Inches(1.5), Inches(11.9), Inches(3.9), size=14,
      col_w=[2.6, 1.3, 1.3, 1.3, 2.9], bold_cols=(4,),
      red_cells=tuple((r, c) for r in (1, 2, 3, 4, 6) for c in range(1, 5)))
bullets(s, [
    ("8/8 targets are only ~600–1,200 donors lower than 10/10 at 95% coverage", {"bold": True, "color": ACCENT}),
    "Because DRB1 and DQB1 versions co-travel ≥94% of the time, a DRB1-matched donor is almost always DQB1-matched too",
    "Practical implication: full 10-gene typing adds clinical value at almost no extra recruitment cost",
], t=Inches(5.55), size=16, gap=6)
notes(s, "Result two — Table 2. What if we use the looser 8-out-of-8 standard, ignoring DQB1 entirely? You might expect the targets to fall a lot. They barely move: typically 600 to 1,200 fewer donors at the 95 per cent level. The reason is the gene-block effect from earlier: DRB1 and DQB1 versions travel together at least 94 per cent of the time, so any donor matched at DRB1 is almost automatically matched at DQB1 too — the extra requirement is nearly free. This has a genuinely useful practical implication: registries can justify full ten-allele typing — which adds real clinical information, because DQB1 mismatches can raise graft-versus-host disease risk — without meaningfully raising their recruitment targets. There is no trade-off here; do the full typing. [~1.5 min]")

# ---------------------------------------------------------------- 15 Table 3 cross-ethnic
s = slide()
title(s, "Result 3 — a shared pool cannot replace same-ethnicity donors")
table(s, [
    ["Patient group", "75%", "85%", "90%", "95% coverage"],
    ["Chinese", "12,198", "25,941", "42,769", "93,348"],
    ["Malay", ">10 million", ">10 million", ">10 million", ">10 million"],
    ["Indian", "1.6 million", "5.2 million", ">10 million", ">10 million"],
    ["Others", ">10 million", ">10 million", ">10 million", ">10 million"],
], Inches(0.7), Inches(1.55), Inches(11.9), Inches(2.9), size=16,
      col_w=[2.2, 1.6, 1.6, 1.6, 1.8],
      red_cells=((2, 1), (2, 2), (2, 3), (2, 4), (3, 1), (3, 2), (3, 3), (3, 4), (4, 1), (4, 2), (4, 3), (4, 4)))
bullets(s, [
    "Model: one combined registry mirroring Singapore's population (74.3% Chinese)",
    ("For Malay, Indian and Others patients, no realistic registry size achieves coverage — their haplotypes simply aren't in the pool", {"bold": True, "color": ACCENT}),
    "Even Chinese patients need ~2× more donors from a mixed pool (93,348 vs 41,183)",
    "Same-ethnicity recruitment is not merely preferable — it is the only viable strategy",
], t=Inches(4.75), size=17, gap=8)
notes(s, "Result three answers the tempting question: couldn't one big shared registry — which in Singapore would be 77 per cent Chinese — serve everyone? Table 3 says no, definitively. For Chinese patients a shared pool works, but inefficiently: you'd need 93,000 donors instead of 43,000, because most of the pool is now diluted from any one patient's perspective. But look at the red rows: for Malay, Indian and Others patients, the model cannot reach 95 per cent coverage with any registry under ten million donors. Their distinctive haplotypes are simply not present in a Chinese-dominated pool — no amount of scale fixes an absence. This mirrors what the Israeli and US registries have reported. The policy conclusion is blunt: same-ethnicity recruitment is not a nice-to-have; for these communities it is the only strategy that works. [~2 min]")

# ---------------------------------------------------------------- 16 Figure 3 partial match 10-locus
s = slide()
title(s, "Result 4 — accepting one mismatch halves the requirement")
pic(s, f"{FIG}/partial_match_10locus.png", Inches(0.7), Inches(1.35), Inches(9.0), Inches(5.3))
bullets(s, [
    ("Green: 10/10 · Blue: ≥9/10 · Red: ≥8/10", {"size": 15}),
    ("Dashed = cross-ethnic pool", {"size": 15}),
    ("Relaxing 10/10 → 9/10 roughly halves N* for every group", {"bold": True, "size": 17}),
    ("Chinese at 95%: ~20–22k (9/10) vs 41,183 (10/10)", {"size": 16}),
    ("Trials show a single permissive mismatch can carry minimal survival impact", {"size": 16}),
], l=Inches(9.85), t=Inches(1.6), w=Inches(3.2), size=16, gap=10)
caption(s, "Figure 3. Coverage vs registry size for 10-gene partial matching, by CMIO group.")
notes(s, "Result four: what does clinical flexibility buy? These are coverage curves — registry size along the bottom, per cent of patients matched on the vertical axis — for each CMIO group. Green is the strict 10-out-of-10 standard; blue allows one mismatch; red allows two. The striking feature is the jump from green to blue: accepting a single mismatched allele roughly halves the required registry for every group. For Chinese patients, 95 per cent coverage needs about 21,000 donors at 9-out-of-10, versus 43,000 at 10-out-of-10. And this isn't a reckless trade: large multi-centre studies show that certain single mismatches — particularly 'permissive' ones — carry minimal additional survival impact. In effect, formalising an evidence-based 9-out-of-10 protocol doubles the reach of a registry without recruiting a single new donor. Note also the dashed cross-ethnic curves flattening out — relaxation doesn't rescue the shared-pool strategy either. [~2 min]")

# ---------------------------------------------------------------- 17 Figure 4 partial match 8-locus
s = slide()
title(s, "The same pattern at the 8-gene standard")
pic(s, f"{FIG}/partial_match_8locus.png", Inches(0.7), Inches(1.35), Inches(9.0), Inches(5.3))
bullets(s, [
    ("Green: 8/8 · Blue: ≥7/8 · Red: ≥6/8", {"size": 15}),
    ("One-allele relaxation gives the same halving effect", {"bold": True, "size": 17}),
    ("Cross-ethnic curves (dashed) plateau for Malay, Indian, Others — partial matching cannot fix a diversity gap", {"size": 16}),
], l=Inches(9.85), t=Inches(1.8), w=Inches(3.2), size=16, gap=12)
caption(s, "Figure 4. Coverage vs registry size for 8-gene partial matching, by CMIO group.")
notes(s, "For completeness, the same analysis in the 8-gene framework — and the same story: relaxing by one allele, from 8-of-8 to 7-of-8, roughly halves the requirement in every group. But I want to draw your eye to the dashed lines again, the cross-ethnic curves. For Malay, Indian and Others patients they plateau well below the coverage targets, no matter how far you relax the matching standard. That's an important negative result: partial matching stretches a registry that already contains roughly the right haplotypes; it cannot conjure up haplotypes that were never recruited. Flexibility and diversity are complements, not substitutes. [~1 min]")

# ---------------------------------------------------------------- 18 sensitivity
s = slide()
title(s, "Result 5 — the targets don't depend on demographic assumptions")
pic(s, f"{FIG}/cross_ethnic_sensitivity.png", Inches(0.55), Inches(1.4), Inches(7.3), Inches(4.4))
table(s, [
    ["Scenario (C/M/I/O weights)", "N* at 95%"],
    ["SG population (74.3/13.5/9/3.2)", "42,567"],
    ["BMDP+SCBB donors (75/9/9/6)", "42,289"],
    ["HSA patient data (72/15/5/8)", "41,950"],
    ["Minority-focus (0/40/40/20)", "41,129"],
], Inches(8.15), Inches(1.7), Inches(4.7), Inches(2.7), size=13, col_w=[3, 1.5])
bullets(s, [
    ("Combined N* varies by 3.5% across all four scenarios", {"bold": True, "size": 16, "color": RED}),
    ("The ~40–45k target is a structural property of CMIO haplotype diversity, robust to demographic change", {"size": 15}),
], l=Inches(8.15), t=Inches(4.6), w=Inches(4.7), size=15, gap=8)
caption(s, "Figure 5 & Table 6. Registry size under four patient-composition scenarios — near-identical bars.")
notes(s, "A planner's natural worry: do these numbers depend on assumptions about who the future patients are? We stress-tested that. Table 6 and Figure 5 show the combined registry size under four very different patient-mix scenarios: Singapore's census composition, the donor registry's own composition, the actual referral mix seen by the Health Sciences Authority, and an extreme minority-only scenario with no Chinese patients at all. The bars are near-identical — the answer moves by less than three per cent. Why? Because every group's individual requirement is in the same 31-to-44-thousand band, reweighting between them barely moves the total. The practical value: the 40-to-45-thousand target is structural. It will not be invalidated by demographic drift, and no one can argue it away by quibbling with the patient-mix assumption. [~1.5 min]")

# ---------------------------------------------------------------- 19 validation
s = slide()
title(s, "Result 6 — checking the model against real patients")
table(s, [
    ["Group", "Shared haplotypes", "Spearman r", "RMSE"],
    ["Chinese", "33", "0.70 (p<0.001)", "0.0094"],
    ["Malay", "11", "n/a*", "0.0284"],
    ["Indian", "1", "n/a*", "—"],
    ["Others", "4", "n/a*", "0.0418"],
], Inches(0.7), Inches(1.5), Inches(6.2), Inches(2.4), size=13, col_w=[1.4, 1.6, 1.7, 1.2])
bullets(s, [
    ("Chinese estimates validate well against 564 real HSA patient–donor pairs", {"bold": True, "size": 16}),
    ("Malay/Indian/Others: too few shared haplotypes → their targets are model projections pending validation", {"size": 15}),
    ("* n < 3 shared haplotypes — rank correlation not meaningful", {"size": 12, "color": MUTED}),
], l=Inches(7.25), t=Inches(1.55), w=Inches(5.5), size=15, gap=8)
pic(s, f"{FIG}/match_validation_scatter.png", Inches(0.7), Inches(4.05), Inches(11.9), Inches(2.7))
caption(s, "Figure 6 & Table 7. Observed patient haplotype frequencies vs model estimates; dashed line = perfect agreement.")
notes(s, "Before trusting a model, check it against reality. We compared our estimated haplotype frequencies with frequencies observed directly in 564 real patient–donor pairs from the Health Sciences Authority. For the Chinese group — where 33 haplotypes overlapped — the agreement is good: a rank correlation of 0.70, highly significant, and small errors. The scatter shows points hugging the perfect-agreement line. For the other groups we simply had too few overlapping haplotypes — eleven for Malay, four for Others, one for Indian — so no meaningful correlation can be computed. We're explicit about this asymmetry in the paper: the Chinese targets are empirically validated; the Malay, Indian and Others targets are model projections pending validation, and growing minority patient cohorts is itself one of our recommendations. [~1.5 min]")

# ---------------------------------------------------------------- 20 Others PCA
s = slide()
title(s, "Result 7 — 'Others' is not one population, it's three")
pic(s, f"{FIG}/others_pca_scatter.png", Inches(0.7), Inches(1.35), Inches(8.6), Inches(5.3))
bullets(s, [
    ("Clustering of 3,847 fully-typed Others donors by HLA profile", {"size": 16}),
    ("Three groups; boundaries overlap (s=0.24)", {"bold": True, "size": 16}),
    ("European / Eurasian", {"size": 15}),
    ("Filipino / SE Asian", {"size": 15}),
    ("Northeast Asian / Mixed", {"size": 15}),
    ("Ancestry inferred from haplotype signatures (AFND)", {"size": 14, "color": MUTED}),
], l=Inches(9.45), t=Inches(1.7), w=Inches(3.6), size=15, gap=8)
caption(s, "Figure 7. PCA of the 3,847 fully-typed Others donors: three ancestry clusters (k=3 optimal by silhouette; s=0.24 in the five-PC clustering space, 0.43 in the projection shown).")
notes(s, "Now the most surprising finding. The 'Others' census category lumps together Eurasians, Europeans, Filipinos, and many mixed backgrounds. Is it statistically legitimate to model them as one population? We let the data answer: we ran principal component analysis and clustering on the HLA profiles of all 3,847 fully-typed Others donors, with no ancestry labels supplied. The result is this picture — three clusters, visually distinct in the PCA projection, with k=3 selected as optimal by the silhouette criterion. Be honest about the strength of that separation: the silhouette is 0.24 in the full clustering space, so the boundaries overlap; what really carries the argument is that each cluster has its own characteristic haplotypes. Matching each cluster's characteristic haplotypes against international reference databases identifies them as European-slash-Eurasian, Filipino-slash-Southeast-Asian, and Northeast Asian or mixed. So 'Others' is not a population — it's at least three, hiding inside one administrative box. And that has direct consequences for the registry target, on the next slide. [~1.5 min]")

# ---------------------------------------------------------------- 21 Others tables
s = slide()
title(s, "Why the pooled 'Others' target is misleading")
table(s, [
    ["Cluster", "Putative ancestry", "N donors", "75%", "85%", "90%", "95%"],
    ["1", "European / Eurasian", "1,029", "3,889", "9,544", "16,845", "35,193"],
    ["2", "Filipino / SE Asian", "1,257", "12,743", "24,836", "37,239", "63,856"],
    ["3", "NE Asian / Mixed", "1,561", "11,681", "20,024", "28,287", "45,731"],
], Inches(0.7), Inches(1.5), Inches(11.9), Inches(2.0), size=14,
      col_w=[0.8, 2.2, 1.1, 1.1, 1.1, 1.1, 1.1], red_cells=((2, 6),))
table(s, [
    ["Cluster", "Top haplotype (A~B~C~DRB1~DQB1)", "Freq", "Signature of"],
    ["1", "A*01:01~B*08:01~C*07:01~DRB1*03:01~DQB1*02:01", "12.2%", "8.1 Ancestral — hallmark N. European"],
    ["2", "A*24:07~B*35:05~C*04:01~DRB1*12:02~DQB1*03:01", "9.1%", "Filipino / SE Asian archipelago"],
    ["3", "A*02:07~B*46:01~C*01:02~DRB1*09:01~DQB1*03:03", "4.7%", "Chinese-specific"],
], Inches(0.7), Inches(3.75), Inches(11.9), Inches(1.9), size=12, col_w=[0.7, 4.6, 0.8, 2.6])
bullets(s, [
    ("Sub-group requirements span 35,193 → 63,856 — the pooled figure of 31,129 is below every one of them", {"bold": True, "color": ACCENT, "size": 16}),
    ("Recommended planning ceiling for Others: 63,856 (the Filipino/SE Asian cluster)", {"bold": True, "size": 16}),
], t=Inches(5.85), size=16, gap=6)
notes(s, "Here's the payoff of that clustering — Tables 4 and 5. Each sub-group has its own registry requirement: 35,000 for the European-Eurasian cluster, 46,000 for the Northeast-Asian-mixed cluster, and 64,000 for the Filipino-Southeast-Asian cluster, which is the most haplotype-diverse. Now recall the pooled 'Others' figure from Table 1: 31,129 — below every single sub-group's true requirement. That's the artefact: mixing three different populations makes the blend look artificially easy to match, because the model wrongly assumes any Others donor can match any Others patient. The lower table shows the evidence for the ancestry labels — each cluster's top haplotype is a textbook population signature, like the 8.1 ancestral haplotype marking Northern Europeans. Planning recommendation: use 63,856 — the highest sub-group ceiling — and start collecting simple ancestry information at donor registration. [~2 min]")

# ---------------------------------------------------------------- 22 honesty / limitations
s = slide()
title(s, "Stress-testing our own assumptions")
pic(s, f"{FIG}/em_convergence.png", Inches(0.55), Inches(1.45), Inches(6.9), Inches(4.6))
bullets(s, [
    ("Computational cap check (Fig. S1): capping the EM at 5,000 samples overestimates the Chinese N* by 8.2% — errs on the safe side", {"size": 15}),
    ("Would a much larger cohort change the answer? No — the haplotype vocabulary above the 0.1% floor saturates (table under the figure), and is 123–144 in all four groups. More donors cannot create common haplotypes, only measure them better", {"size": 14, "color": RED}),
    ("Key limitations:", {"bold": True, "size": 16}),
    (1, "0.1% frequency floor \u2192 all coverage conditional; at a 1\u00d710\u207b\u2076 floor the Chinese 95% target becomes 8.7\u00d710\u2077 \u2014 absolute sizes are order-of-magnitude", {"size": 14}),
    (1, "CIs cover frequency-sampling noise only", {"size": 14}),
    (1, "HWE departures → Indian & Others estimates exploratory", {"size": 14}),
    (1, "Minority validation limited by small patient cohorts", {"size": 14}),
    (1, "Rare-haplotype smoothing shifts 95% targets by ≤3%", {"size": 14}),
    (1, "Real-world attrition ~30–50% → sign-up targets ≈ 1.67 × N*", {"size": 14}),
], l=Inches(7.75), t=Inches(1.5), w=Inches(5.1), size=15, gap=6)
table(s, [
    ["Chinese donors used", "Haplotypes ≥0.1%", "N* at 95%"],
    ["5,000  (the cap)", "143", "45,148"],
    ["20,000", "140", "44,326"],
    ["45,018  (all)", "136", "41,727"],
], Inches(0.55), Inches(5.80), Inches(6.90), Inches(0.98), size=12,
      col_w=[2.4, 2.2, 1.8],
      red_cells=tuple((r, c) for r in range(1, 4) for c in range(3)))
caption(s, "Figure S1. Chinese N* vs EM sample size: stable above ~20,000 samples; the 5,000 cap (red dashed) is conservative. "
           "Table: the vocabulary above the floor saturates — 9× more donors changes it by 7 haplotypes.")
notes(s, "Every model has assumptions, and we stress-tested ours. The figure addresses the most technical one: for computational reasons the haplotype estimation was capped at 5,000 donors per group. Only the Chinese group has more data than that, so we re-ran the estimation at sample sizes from 500 up to the full 45,000. The curve stabilises above 20,000 samples, and at the cap the estimate is 8.2 per cent too high — an error in the safe direction: we'd recruit slightly more donors than strictly needed, never fewer. The same run answers a question people often ask: would a much larger cohort — say half a million donors — change these targets? Look at how many haplotypes clear the 0.1 per cent floor: 143 at five thousand donors, 136 at forty-five thousand. It saturates. At a fixed floor, more donors cannot make rare haplotypes common; they only measure the common ones more precisely. And the flip side is the important one — the mass sitting below the floor is a property of the population, not of our sample size, so no amount of recruitment lifts those patients above it. That is why mismatch tolerance, not scale, is the lever for them. Other limitations, briefly: the 0.1 per cent frequency floor makes every coverage figure conditional, and moving the floor moves the absolute targets by orders of magnitude — the comparative findings are what survive; the confidence intervals cover sampling noise only; Hardy–Weinberg departures make the Indian and Others estimates more exploratory; minority validation awaits larger patient cohorts; rare-haplotype smoothing moves the 95 per cent targets by under three per cent; and real-world attrition means sign-up targets exceed the biological targets by roughly two-thirds. [~1.5 min]")

# ---------------------------------------------------------------- 23 recommendations
s = slide()
title(s, "Six recommendations")
bullets(s, [
    ("1. Build same-ethnicity registries of 40,000–45,000 donors per major CMIO group — this buys ~49% full-match coverage, not 95%", {"bold": True}),
    ("2. Prioritise Malay and Indian recruitment urgently — each needs ~40,000, comparable to today's entire BMDP", {"bold": True}),
    ("3. Formalise evidence-based 9/10 partial-match protocols — a 9/10 donor is an MMUD; doubles reach at zero cost and is the only route for below-floor patients", {"bold": True}),
    "4. Plan Others recruitment to the 63,856 ceiling and collect ancestry data at registration",
    "5. Keep full DQB1 typing — clinical value at negligible extra registry cost",
    "6. Keep 95% (conditional) coverage as the planning standard — but neither 90% nor 95% reaches patients below the frequency floor; their route is Recommendation 3",
], w=Inches(9.4), size=19, gap=11)
pic(s, "assets/cartoons/checklist.png", Inches(10.2), Inches(1.7), Inches(2.8), Inches(4.7))
notes(s, "Pulling it all together, six recommendations. One: build same-ethnicity registries of 40 to 45 thousand donors per major group — the model shows there is no shared-registry shortcut. Two: Malay and Indian recruitment is the urgent gap; each community needs a registry comparable in size to today's entire BMDP. Three: formalise 9-out-of-10 partial-match protocols where the evidence supports them — a 9-out-of-10 donor is by definition a mismatched unrelated donor, so this is the MMUD recommendation: it doubles reach for free, and for patients whose haplotypes fall below the frequency floor it is the only route to a donor at all. The safety evidence comes from the external trial literature, not from this analysis. Four: for the Others group, plan to the 64,000 ceiling and start asking two simple ancestry questions at registration, so the pooled figure can be retired. Five: keep full DQB1 typing — the linkage structure makes it clinically valuable and essentially free. Six: keep 95 per cent as the planning standard — but be honest that both 90 and 95 are conditional numbers. Patients below the frequency floor are not reached at either threshold; their access runs through recommendation three. [~2 min]")

# ---------------------------------------------------------------- 24 conclusions
s = slide()
title(s, "Take-home messages")
bullets(s, [
    ("≈ 40,000–45,000 same-ethnicity donors per CMIO group gives 95% of ELIGIBLE patients a full match — about 49% of all patients; 41.8% observed today", {"bold": True, "size": 20}),
    ("Cross-ethnic matching cannot substitute — for Malay, Indian and Others patients, no realistic shared registry works", {"size": 21}),
    ("Accepting one mismatch (9/10) halves the requirement — the cheapest expansion available", {"size": 21}),
    ("'Others' hides three distinct populations — plan to the 63,856 ceiling", {"size": 21}),
    ("First quantitative, uncertainty-quantified registry targets for Singapore — a template for multiethnic registries across Asia-Pacific", {"size": 21, "color": ACCENT, "bold": True}),
], w=Inches(8.9), size=21, gap=14)
pic(s, "assets/cartoons/target.png", Inches(9.7), Inches(2.2), Inches(3.2), Inches(3.2))
notes(s, "So, four take-home messages. First, the number: about 40 to 45 thousand same-ethnicity donors per community gives 95 per cent of common-haplotype patients a fully matched donor — which, once the frequency floor is accounted for, is roughly half of all patients, and 41.8 per cent is what the Chinese registry actually delivers today. Second, there is no shortcut through a shared pool: for Malay, Indian and Others patients, cross-ethnic matching fails at any realistic scale. Third, clinical flexibility is powerful: accepting a single well-chosen mismatch halves the requirement overnight. And fourth, the Others category conceals three genetically distinct populations and must be planned to its hardest sub-group, not its misleading pooled average. More broadly, this is the first uncertainty-quantified registry sizing for Singapore, and the framework transfers directly to other multiethnic populations across the Asia-Pacific. [~1.5 min]")

# ---------------------------------------------------------------- 25 thanks
s = slide()
txt(s, Inches(1.0), Inches(2.5), Inches(11.3), Inches(1.2), "Thank you", size=44, bold=True, color=ACCENT)
txt(s, Inches(1.0), Inches(3.9), Inches(11.3), Inches(1.8),
    "Data: Bone Marrow Donor Programme · Singapore Cord Blood Bank · Health Sciences Authority\n"
    "Analysis code and data: github.com/alvin8-git/HLA\n"
    "Contact: alvin1976sg@gmail.com", size=18, color=INK)
txt(s, Inches(1.0), Inches(5.9), Inches(11.3), Inches(0.6),
    "Ng AYJ et al., Blood Cell Therapy 2022 (foundation dataset) · Manuscript v2.15c", size=14, color=MUTED)
pic(s, "assets/cartoons/skyline.png", Inches(6.9), Inches(5.35), Inches(6.0), Inches(1.9))
notes(s, "Thank you. I'd like to acknowledge the Bone Marrow Donor Programme, the Singapore Cord Blood Bank, and the Health Sciences Authority for the data that made this possible. All analysis code and intermediate data are openly available on GitHub — the entire pipeline is reproducible. I'm happy to take questions — on the method, the numbers, or what it would take to actually recruit 40,000 Malay donors. [~0.5 min, then Q&A]")

# ── Reader notes: one "Additional detail" paragraph appended per slide ────────
# Keyed by 1-based slide number. Speaking script stays first; this paragraph is
# for a reader working through the deck without the talk.
EXTRA = {
    1: "Additional detail — Deck accompanies manuscript v2.15c (registry-size "
       "model, 0.1% haplotype frequency floor, frozen 1e-3 snapshot data). All "
       "registry sizes are conditional on the patient carrying two haplotypes "
       "above the floor; see slide 12 and manuscript Section 4.1.",
    2: "Additional detail — ~70% of HSCT candidates lack a matched sibling and "
       "depend on unrelated donors. Matching quality is the strongest "
       "registry-controllable predictor of transplant outcome (manuscript "
       "Section 1, ref [7]).",
    3: "Additional detail — Ng et al., Blood Cell Therapy 2022 characterised "
       "allele and haplotype frequencies for the same 59,186-donor cohort; the "
       "present study converts those frequencies into registry-size targets "
       "(manuscript Section 1, ref [1]).",
    4: "Additional detail — Matching standards: 10/10 = HLA-A, -B, -C, -DRB1, "
       "-DQB1 at two-field resolution; 8/8 omits DQB1. DPB1 is not modelled "
       "(not typed in the source data) — a stated limitation, as permissive "
       "DPB1 mismatching is standard-of-care in Singapore HSCT.",
    5: "Additional detail — The four questions map to manuscript Sections 3.1 "
       "(same-ethnicity targets), 3.3 (cross-ethnic), 3.4 (partial match), and "
       "3.7 (Others stratification). Coverage targets are conditional; see "
       "slide 12.",
    6: "Additional detail — Counts by group (5-locus complete): Chinese 44,400, "
       "Malay 5,578, Indian 5,490, Others 3,767 = 59,235 re-extracted; 59,186 "
       "after the published inclusion criteria (manuscript Section 2.1). HSA "
       "records (1,350 donors + 564 patients) are held out for validation and "
       "the demographic scenario only.",
    7: "Additional detail — stage by stage through the flowchart. Each stage "
       "is given first in plain terms (why we do it at all), then in "
       "technical terms.\n"
       "INPUT DATA (blue). Why: before you can ask how many donors a country "
       "needs, you need an honest picture of what tissue types Singaporeans "
       "actually carry — and that picture has to be built separately for each "
       "community, because the common types differ between them. A lab test "
       "tells us a donor carries, say, two versions of gene A, but not how "
       "those versions are grouped together on their two chromosomes — like "
       "being told a person owns a red and a blue shoe of each kind, without "
       "being told which pair they actually wear. How: 59,186 registry donors "
       "typed at the five HLA genes (A, B, C, DRB1, DQB1), cleaned to one row "
       "per person per gene and split into the four CMIO groups.\n"
       "EM HAPLOTYPE PHASING (green, Key idea 1). Why: the five genes sit "
       "side by side on one chromosome and are inherited as a block, so "
       "certain combinations are far more common than chance would suggest. "
       "If you ignore that and treat the genes as independent, you get the "
       "arithmetic badly wrong — common patients look rarer than they are, "
       "and the registry you calculate is the wrong size in the wrong places. "
       "So we have to work out the blocks, not just the individual genes. "
       "How: the Excoffier–Slatkin expectation–maximisation algorithm infers "
       "the five-gene haplotypes and their frequencies per group; run on up "
       "to 5,000 people per group for speed (a known, conservative "
       "approximation) and validated against the international GENE[RATE] "
       "reference pipeline.\n"
       "0.1% FREQUENCY FLOOR. Why: in a sample of this size, a type seen once "
       "or twice could genuinely be that rare, or could be a typing error, or "
       "could be pure luck of who walked through the door. Building "
       "recruitment targets on those numbers would be building on noise, so "
       "we set them aside and are explicit that we have done so — the honest "
       "position is a firm answer for common patients plus a stated blind "
       "spot, rather than a shaky answer for everybody. How: haplotypes below "
       "1-in-1,000 are dropped and the rest renormalised; this is why every "
       "result is CONDITIONAL on the patient's two haplotypes both being "
       "common (slide 12 quantifies who that leaves out).\n"
       "HWE EXPANSION (green). Why: patients aren't matched one block at a "
       "time — a transplant match means the donor matches the patient's whole "
       "tissue type, both blocks, all ten alleles. So we need the frequency "
       "of complete types, not of blocks: two individually common blocks can "
       "make an uncommon combination, and it's the combination the patient "
       "has to find. How: haplotypes are paired under random-mating "
       "(Hardy–Weinberg) assumptions to give the frequency F of each full "
       "ten-allele genotype.\n"
       "COVERAGE MODEL (green, Key idea 2). Why: this is the question the "
       "whole paper exists to answer — if the registry held N donors, what "
       "fraction of patients would find their match? Common patients are "
       "found almost immediately; rare ones need a much bigger pool. "
       "Averaging across all patients, weighted by how common each type is, "
       "turns 'how big is the registry' into 'how many patients go home with "
       "a donor'. How: a patient with genotype frequency F matches with "
       "probability 1−(1−F)^N, and population coverage is "
       "C(N) = Σ F·[1−(1−F)^N].\n"
       "REGISTRY-SIZE SEARCH (yellow). Why: recruiting donors costs real "
       "money and real volunteers, so the useful output is not a curve but a "
       "number — the smallest registry that reaches the standard we set "
       "ourselves. We run the coverage calculation backwards: pick the target "
       "first, then ask what size delivers it. How: C(N) rises monotonically "
       "with N, so a log-scale binary search returns the smallest N* meeting "
       "each target (75–95%). Near 95% the answer is driven by the rarest "
       "patients still above the floor — which is why N* is so sensitive to "
       "where that floor sits.\n"
       "BOOTSTRAP (red, Key idea 3). Why: every frequency here was measured "
       "from one particular set of 59,186 people. Recruit a different 59,186 "
       "Singaporeans and the numbers would come out slightly differently — so "
       "quoting a single figure would imply a precision we don't have. We "
       "therefore ask how much the answer would wobble, and publish the "
       "wobble alongside the estimate. How: haplotype frequencies are redrawn "
       "1,000 times from a Dirichlet distribution scaled to the sample size, "
       "the whole calculation is rerun each time, and we report the median "
       "with a 2.5th–97.5th percentile range. Important limit: these bars "
       "cover sampling noise only — they do not cover the modelling "
       "assumptions above, which the limitations section handles.",
    8: "Additional detail — DRB1–DQB1 D' = 0.94–0.99 in all four groups; B–C "
       "D' ≥ 0.95. This LD is why 8/8 targets sit only 600–1,200 donors below "
       "10/10 (manuscript Sections 2.2 and 3.2, Table 2).",
    9: "Additional detail — Match probability for one patient: 1−(1−F)^N. "
       "Population coverage weights this by diplotype frequency. The model "
       "assumes independent draws and registry composition mirroring the donor "
       "pool (manuscript Section 2.3, equations 1–4).",
    10: "Additional detail — Dirichlet parametric bootstrap, B=1,000; reported "
        "point estimate is the bootstrap median, range = 2.5th–97.5th "
        "percentile. Concentration is n_eff × f, with n_eff the number of "
        "individuals the frequencies were estimated from — the 5-locus donor "
        "count subject to the 5,000-person EM limit (Chinese/Malay/Indian "
        "5,000; Others 3,767), not the full group size. Using the full size "
        "for Chinese (44,400) would narrow that interval about fourfold and "
        "claim a precision the frequencies do not carry. The manuscript "
        "(Section 2.4) notes these cover frequency-sampling noise only.",
    11: "Additional detail — Full table: manuscript Table 1. 95% conditional "
        "targets: Chinese 41,183 (40,184–42,153), Malay 39,831 "
        "(38,680–40,935), Indian 43,785 (42,762–44,547), Others pooled 31,129 "
        "(30,420–32,001; artefact — see slide 22). Signed-up targets gross up "
        "by ÷0.60 for 40% attrition.",
    12: "Additional detail — Retained frequency mass: Chinese 51.7%, Malay "
        "52.9%, Indian 40.2%, Others 35.6% (manuscript Sections 2.2 and 4.1). "
        "Observed twin rates are exact-genotype leave-one-out counts — no EM, "
        "no HWE, no floor. 0.95 × 0.517 = 49.1% is the unconditional Chinese "
        "equivalent; 41.8% is observed; the gap direction is consistent with "
        "population substructure (Wahlund effect).",
    13: "Additional detail — Manuscript Figure 2. Ranges reflect sampling "
        "variability within the retained haplotype set; they exclude EM, HWE, "
        "and floor uncertainty, and the floor dominates (manuscript Section "
        "2.4, third limitation).",
    14: "Additional detail — The long tail is why coverage grows "
        "logarithmically: each doubling of the registry adds a thinner slice "
        "of rare diplotypes. Manuscript Figure (diplotype long-tail) and "
        "Section 3.1.",
    15: "Additional detail — Manuscript Table 2. 8/8 targets at 95%: Chinese "
        "42,115, Malay 36,176, Indian 42,738, Others 30,525. The small "
        "10/10→8/8 saving is the DRB1–DQB1 LD of slide 8.",
    16: "Additional detail — Manuscript Table 3 and Section 3.3. At Singapore "
        "weights (74.3% Chinese), a pooled registry reaches Chinese patients at "
        "~93k but cannot reach 95% for Malay, Indian, or Others at any size "
        "examined (ceiling 1×10⁷ in the v2.15 run) — their haplotypes are "
        "absent from the majority pool.",
    17: "Additional detail — Manuscript Section 3.4, Figure 3. 9/10 halves N* "
        "for every group (Chinese 95%: ~21k vs 41,183). A 9/10 unrelated donor "
        "is by definition an MMUD; safety at permissive mismatches rests on "
        "external trials (refs [12,15]) — this analysis contains no outcome "
        "data.",
    18: "Additional detail — Manuscript Figure 4 (8-locus framework: 8/8, 7/8, "
        "6/8). Same halving pattern; cross-ethnic curves stay infeasible at "
        "every relaxation level, so relaxation does not rescue a shared-pool "
        "strategy.",
    19: "Additional detail — Manuscript Table 6, Section 3.5. Four scenarios "
        "(SG population 74.3/13.5/9/3.2, registry composition 75/9/9/6, HSA referral "
        "72/15/5/8, minority-focus 0/40/40/20) move the combined N* by 3.5%. "
        "The 2020 census citizens+PR breakdown (73.3/13.5/9.0/3.2) lies within "
        "this envelope.",
    20: "Additional detail — Manuscript Table 7, Section 3.6. Chinese: "
        "Spearman r = 0.70 on shared haplotypes vs 564 HSA patient–donor "
        "pairs. Malay/Indian/Others have too few shared haplotypes (n < 3) — "
        "their targets are model projections pending validation.",
    21: "Additional detail — PCA on binary allele indicators (alleles ≥1%), "
        "k-means, optimal k=3 by silhouette; s = 0.24 in the five-PC "
        "clustering space (0.43 in the PC1–PC2 projection shown) — clusters "
        "overlap; read as a caution against pooling, not as established "
        "sub-populations (manuscript Sections 2.6, 3.7, 4.1).",
    22: "Additional detail — Manuscript Tables 4–5. Cluster targets at 95%: "
        "European/Eurasian 35,193, NE Asian/Mixed 45,731, Filipino/SE Asian "
        "63,856. The pooled 31,129 sits below all three because mixing "
        "heterogeneous pools inflates apparent common-haplotype coverage.",
    23: "Additional detail — Would a much larger cohort change these numbers? "
        "No, and the convergence run shows why. The count of haplotypes "
        "clearing the 0.1% floor is 143 at n=5,000, 140 at n=20,000 and 136 at "
        "the full n=45,018 — it saturates — and sits at 123–144 across all four "
        "groups despite a 12-fold span in group size (Others 3,767 to Chinese "
        "44,400). At a fixed floor extra donors cannot create common "
        "haplotypes; they only estimate the existing ones more precisely, which "
        "is why N* plateaus by n≈5,000 and moves only ~6% from 20,000 to "
        "45,018. The n=500 figure of 451 haplotypes is the opposite artefact: "
        "at that depth almost everything observed appears to clear the floor.\n"
        "The corollary matters more. Retained mass (51.7/52.9/40.2/35.6%) "
        "reflects population structure, not sample size, so a much larger "
        "registry would not lift a single below-floor patient above it — "
        "recruitment scale cannot close that gap, mismatch tolerance "
        "(Recommendation 3) can. What a larger cohort would buy is a defensible "
        "floor below 0.1%, since 1/(2n) only reaches ~10⁻⁶ at that scale, plus "
        "the removal of the small-group caveats in Sections 3.6 and 4.1.\n"
        "Floor sensitivity (manuscript Sections 2.2 and 4.1): a 1×10⁻⁶ floor "
        "(below every group's singleton frequency 1/(2n)) raises the Chinese "
        "95% target from 41,183 to 8.7×10⁷. Absolute sizes are "
        "order-of-magnitude guidance; ratios computed under a common floor "
        "(cross-ethnic penalty, relaxation gains) are the durable results.",
    24: "Additional detail — per recommendation (all map to manuscript "
        "Section 4):\n"
        "R1 (40,000–45,000 per group): Table 1 bootstrap medians at 95% "
        "conditional 10/10 coverage — Chinese 41,183, Malay 39,831, Indian "
        "43,785. Conditional on the 0.1% floor: unconditionally these serve "
        "~49/50/38/34% of all patients (slide 12), and 41.8% is observed for "
        "Chinese. Signed-up recruitment must gross up ×1.67 for ~40% donor "
        "attrition (Tables 1–2, bottom rows). Absolute sizes are "
        "order-of-magnitude guidance (Section 4.1).\n"
        "R2 (Malay/Indian urgency): justified by each group's own gap, not by "
        "cross-group ranking (which Section 4.1 rules unsafe): registries of "
        "5,578 and 5,490 stand against ~40–44k targets — a 7–8× shortfall — "
        "and observed 10/10 twin rates are 25.3% and 19.7% vs Chinese 41.8%, "
        "a model-free measure of the same gap.\n"
        "R3 (9/10 protocols = MMUD): relaxing one allele roughly halves N* in "
        "every group (Chinese 95%: ~21k vs 41,183; Section 3.4, Figure 3). A "
        "9/10 unrelated donor is by definition an MMUD; for patients below "
        "the frequency floor this is the only route to a donor at all. "
        "Clinical safety (permissive DPB1, PTCy-era MMUD) rests on external "
        "trials [12,15] — this paper contains no outcome data.\n"
        "R4 (Others ceiling 63,856 + ancestry data): the pooled 31,129 is an "
        "artefact sitting below every sub-cluster (35,193 / 45,731 / 63,856; "
        "Section 3.7). Plan to the hardest cluster; clusters are exploratory "
        "(silhouette 0.24), so an ancestry field at registration is what "
        "eventually replaces inference with observation.\n"
        "R5 (keep DQB1 typing): 8/8 targets sit only 600–1,200 donors below "
        "10/10 because DRB1–DQB1 D' = 0.94–0.99 (Section 3.2, Table 2) — so "
        "full typing costs almost nothing in registry size while DQB1 "
        "mismatches carry GvHD risk [12].\n"
        "R6 (95% standard, honestly stated): 90%→95% conditional roughly "
        "doubles the registry (~23k→~43k). Both thresholds are conditional: "
        "patients below the floor are reached at neither, so the equity "
        "argument runs through R3, not through the threshold choice.",
    25: "Additional detail — The honest summary chain: 95% conditional → "
        "×0.517 mass → 49.1% of all Chinese patients → 41.8% observed. "
        "Recruitment stays worthwhile (10× registry ≈ +20 points at 10/10) "
        "but cannot close the gap; mismatch tolerance can.",
    26: "Additional detail — Reproducibility: build_report_v215c.py + "
        "analysis/snapshot_1e-3/ rebuild the manuscript; "
        "build_slides_v215c.py rebuilds this deck. Additions over v2.15 are "
        "marked red in the manuscript (5 paragraphs).",
}
for idx, sl in enumerate(prs.slides, 1):
    extra = EXTRA.get(idx)
    if not extra:
        continue
    tf = sl.notes_slide.notes_text_frame
    para = tf.add_paragraph()   # blank separator
    para = tf.add_paragraph()
    para.text = extra

prs.save(OUT)
print(f"Saved {OUT} with {len(prs.slides._sldIdLst)} slides")
